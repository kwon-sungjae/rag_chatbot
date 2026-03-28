from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as qn
from lxml import etree
import copy

# ─── Color Palette ─────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x3E)   # 제목 배경
BLUE    = RGBColor(0x1A, 0x56, 0xDB)   # 강조 파랑
LIGHT   = RGBColor(0xEF, 0xF4, 0xFF)   # 연한 배경
ACCENT  = RGBColor(0xF5, 0xA6, 0x23)   # 오렌지 포인트
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1E, 0x1E, 0x2E)
GRAY    = RGBColor(0x64, 0x74, 0x8B)
GREEN   = RGBColor(0x05, 0x96, 0x69)
RED     = RGBColor(0xDC, 0x26, 0x26)

SLD_W = Inches(13.33)
SLD_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLD_W
prs.slide_height = SLD_H

BLANK = prs.slide_layouts[6]   # 완전 빈 레이아웃


# ─── Low-level helpers ─────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text, font_size=Pt(14), bold=False,
                color=DARK, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=Pt(13), bold=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False, level=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


# ─── Slide-layout helpers ──────────────────────────────────────
def slide_header(slide, title, subtitle=None, progress=None, total=15):
    """상단 네이비 헤더 바 + 제목"""
    add_rect(slide, 0, 0, SLD_W, Inches(1.15), fill_rgb=NAVY)
    add_textbox(slide, Inches(0.4), Inches(0.12), Inches(11), Inches(0.7),
                title, font_size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.75), Inches(10), Inches(0.38),
                    subtitle, font_size=Pt(13), color=RGBColor(0xB0, 0xC4, 0xDE),
                    align=PP_ALIGN.LEFT)
    # 슬라이드 번호
    if progress:
        add_textbox(slide, Inches(12.2), Inches(0.35), Inches(1.0), Inches(0.45),
                    progress, font_size=Pt(12), color=RGBColor(0xB0, 0xC4, 0xDE),
                    align=PP_ALIGN.RIGHT)
    # 하단 라인
    add_rect(slide, 0, SLD_H - Inches(0.35), SLD_W, Inches(0.35), fill_rgb=LIGHT)
    add_textbox(slide, Inches(0.3), SLD_H - Inches(0.32), Inches(6), Inches(0.28),
                "RAG를 이용한 효율적인 챗봇 시스템  |  권성재  |  인공지능학과",
                font_size=Pt(9), color=GRAY)


def bullet_box(slide, l, t, w, h, title, items, title_color=BLUE,
               item_color=DARK, item_size=Pt(13), title_size=Pt(15),
               bg_color=None, border_color=BLUE):
    """둥근 카드 형태 bullet box"""
    box = add_rect(slide, l, t, w, h, fill_rgb=bg_color or WHITE,
                   line_rgb=border_color, line_width=Pt(1.5))
    # 제목 바
    add_rect(slide, l, t, w, Inches(0.42), fill_rgb=title_color)
    add_textbox(slide, l + Inches(0.15), t + Inches(0.04), w - Inches(0.3), Inches(0.38),
                title, font_size=title_size, bold=True, color=WHITE)
    # 내용 텍스트박스
    txb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.5),
                                   w - Inches(0.3), h - Inches(0.6))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = item
        run.font.size = item_size
        run.font.color.rgb = item_color
    return box


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)

# 배경 그라데이션 느낌 (상단 네이비 블록)
add_rect(sld, 0, 0, SLD_W, Inches(4.2), fill_rgb=NAVY)
add_rect(sld, 0, Inches(4.2), SLD_W, Inches(3.3), fill_rgb=WHITE)

# 포인트 라인
add_rect(sld, 0, Inches(4.0), SLD_W, Inches(0.18), fill_rgb=ACCENT)

# 영문 부제목 배경
add_rect(sld, Inches(1.2), Inches(0.55), Inches(10.9), Inches(0.55),
         fill_rgb=RGBColor(0x1A, 0x2E, 0x6E))

add_textbox(sld, Inches(1.2), Inches(0.57), Inches(10.9), Inches(0.48),
            "Master's Thesis Proposal  |  Department of Artificial Intelligence",
            font_size=Pt(13), color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)

# 메인 제목
add_textbox(sld, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.0),
            "RAG를 이용한 효율적인 챗봇 시스템",
            font_size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sld, Inches(0.8), Inches(2.35), Inches(11.7), Inches(0.7),
            "An Efficient Chatbot System Using Retrieval-Augmented Generation (RAG)",
            font_size=Pt(18), color=RGBColor(0xB0, 0xD4, 0xFF), align=PP_ALIGN.CENTER,
            italic=True)

# 키워드 태그
kw_x = Inches(3.5)
for kw in ["RAG", "Query Rewrite", "Hybrid Search", "LLM"]:
    add_rect(sld, kw_x, Inches(3.3), Inches(1.6), Inches(0.42),
             fill_rgb=BLUE, line_rgb=None)
    add_textbox(sld, kw_x + Inches(0.05), Inches(3.32), Inches(1.5), Inches(0.38),
                kw, font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    kw_x += Inches(1.75)

# 발표자 정보
info_items = [
    ("발표자",  "권성재  (학번: 4120411004)"),
    ("지도교수", "노동건 교수님"),
    ("소속",    "인공지능학과"),
    ("날짜",    "2026년 3월"),
]
iy = Inches(4.55)
for label, val in info_items:
    add_textbox(sld, Inches(3.5), iy, Inches(1.4), Inches(0.38),
                label, font_size=Pt(13), bold=True, color=BLUE)
    add_textbox(sld, Inches(4.9), iy, Inches(5.0), Inches(0.38),
                val, font_size=Pt(13), color=DARK)
    iy += Inches(0.48)

# 구분선
add_rect(sld, Inches(3.4), Inches(4.45), Inches(6.5), Pt(1.5), fill_rgb=BLUE)
add_rect(sld, Inches(3.4), Inches(6.85), Inches(6.5), Pt(1.5), fill_rgb=ACCENT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — 연구 배경 및 동기 (1): 기업 환경의 정보 접근 문제
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "연구 배경 및 동기", "기업 환경에서의 정보 접근 문제", "01 / 14")

# 세 개의 문제 카드
card_w = Inches(3.9)
card_h = Inches(3.5)
card_y = Inches(1.4)
cards = [
    ("📂 분산된 기업 정보",
     ["• 공지사항·규정·매뉴얼 등\n  방대한 문서가 여러 시스템에 분산",
      "• 직원이 필요 정보를\n  신속하게 찾기 어려움",
      "• 정보 검색에 과도한\n  시간·비용 낭비"]),
    ("🤖 단순 챗봇의 한계",
     ["• 규칙 기반 챗봇은\n  정형화된 질문만 처리 가능",
      "• 자연어 질의에\n  유연한 대응 불가",
      "• 최신 정보 반영 및\n  맥락 이해 불가"]),
    ("💡 LLM 기반 챗봇의 등장",
     ["• GPT 계열 LLM이\n  자연어 이해·생성 혁신",
      "• 그러나 학습 데이터 외\n  내부 문서에 대한 환각(hallucination) 문제",
      "• 기업 전용 지식 기반\n  연동 기술 필요"]),
]
cx = Inches(0.35)
colors = [BLUE, RGBColor(0x6B, 0x21, 0xA8), GREEN]
for (title, bullets), color in zip(cards, colors):
    bullet_box(sld, cx, card_y, card_w, card_h, title, bullets,
               title_color=color, border_color=color, item_size=Pt(12.5))
    cx += card_w + Inches(0.25)

# 하단 강조 문장
add_rect(sld, Inches(0.35), Inches(5.1), Inches(12.6), Inches(0.62),
         fill_rgb=RGBColor(0xFE, 0xF9, 0xEE), line_rgb=ACCENT, line_width=Pt(1.5))
add_textbox(sld, Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.5),
            "▶  기업 내부 문서를 신뢰할 수 있는 근거로 활용하여 정확하게 답변하는 "
            "지능형 챗봇 시스템이 요구됨",
            font_size=Pt(13.5), bold=True, color=RGBColor(0x92, 0x40, 0x00))


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — 연구 배경 및 동기 (2): RAG 기술 현황과 한계
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "연구 배경 및 동기", "RAG 기술 현황 및 기존 접근의 한계", "02 / 14")

# 왼쪽: RAG 개요
lw = Inches(6.1)
bullet_box(sld, Inches(0.35), Inches(1.35), lw, Inches(4.05),
           "🔍  RAG (Retrieval-Augmented Generation) 개요",
           ["• Lewis et al. (2020) 에서 제안된 검색-증강 생성 패러다임",
            "• 외부 지식 베이스에서 관련 문서를 검색(Retrieve)하여\n"
            "  LLM의 생성(Generate)에 맥락(Context)으로 제공",
            "• 학습 데이터 외 최신·도메인 특화 정보 활용 가능",
            "• 환각(Hallucination) 문제를 문서 기반 답변으로 완화",
            "• 기업 챗봇·QA 시스템에 광범위하게 적용"],
           title_color=BLUE, item_size=Pt(13))

# 오른쪽: 한계
rw = Inches(6.5)
bullet_box(sld, Inches(6.6), Inches(1.35), rw, Inches(4.05),
           "⚠️  기존 RAG의 주요 한계",
           ["• [검색 품질] 단순 키워드·벡터 검색만으로는\n"
            "  의미적으로 유사한 문서 누락 가능",
            "• [쿼리 불명확성] 구어체·오탈자·모호한 질의에\n"
            "  벡터 검색 성능 저하",
            "• [희소 vs. 밀집] BM25(희소) 와 벡터(밀집) 검색이\n"
            "  각각의 장단점을 가지며 단독 사용 시 성능 한계",
            "• [평가 미흡] 검색 단계와 생성 단계를 분리한\n"
            "  체계적 평가 연구 부족"],
           title_color=RED, border_color=RED, item_size=Pt(13))

# 화살표(→) 연결 힌트
add_textbox(sld, Inches(6.0), Inches(3.1), Inches(0.65), Inches(0.5),
            "→", font_size=Pt(28), bold=True, color=ACCENT)

# 하단 인용
add_textbox(sld, Inches(0.35), Inches(5.55), Inches(12.6), Inches(0.35),
            "참고: Lewis et al. (2020), Gao et al. (2024 RAG Survey), Ma et al. (2023 Query Rewrite), Chen et al. (2022 Hybrid Retrieval)",
            font_size=Pt(10), color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — 문제 정의 (1): 핵심 문제와 연구 질문
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "문제 정의", "핵심 문제 및 연구 질문", "03 / 14")

# 문제 정의 3-박스
prob_items = [
    ("P1", "쿼리 불명확성 문제",
     "사용자의 구어체·약어·오탈자 질의가 임베딩 공간에서\n"
     "올바른 문서 벡터와 멀어져 관련 문서를 놓치는 현상",
     RGBColor(0xDC, 0x26, 0x26)),
    ("P2", "단일 검색 방법의 한계",
     "벡터 검색은 의미적 유사도 우수, 정확 키워드 매칭 취약.\n"
     "BM25는 키워드 강점, 의미 파악 취약 — 각각 단독 사용 시 성능 천장 존재",
     RGBColor(0xD9, 0x7A, 0x06)),
    ("P3", "검색-생성 연계 평가 부재",
     "검색 성능(Recall@k, MRR)과 생성 품질(EM, F1)을 통합적으로 분석한\n"
     "실험 연구가 부족하여 최적 전략 선택 기준이 불명확",
     RGBColor(0x05, 0x66, 0x99)),
]
by = Inches(1.42)
for pid, ptitle, pdesc, color in prob_items:
    # 번호 원형 배경
    add_rect(sld, Inches(0.35), by, Inches(0.7), Inches(1.18), fill_rgb=color)
    add_textbox(sld, Inches(0.35), by + Inches(0.18), Inches(0.7), Inches(0.7),
                pid, font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 내용 카드
    add_rect(sld, Inches(1.1), by, Inches(11.85), Inches(1.18),
             fill_rgb=WHITE, line_rgb=color, line_width=Pt(1.5))
    add_textbox(sld, Inches(1.25), by + Inches(0.06), Inches(4.5), Inches(0.42),
                ptitle, font_size=Pt(15), bold=True, color=color)
    add_textbox(sld, Inches(1.25), by + Inches(0.5), Inches(11.4), Inches(0.62),
                pdesc, font_size=Pt(12.5), color=DARK)
    by += Inches(1.32)

# 연구 질문
add_rect(sld, Inches(0.35), Inches(5.42), Inches(12.6), Inches(1.65),
         fill_rgb=RGBColor(0xEF, 0xF4, 0xFF), line_rgb=BLUE, line_width=Pt(1.5))
add_textbox(sld, Inches(0.55), Inches(5.5), Inches(3.0), Inches(0.38),
            "핵심 연구 질문 (RQ)", font_size=Pt(14), bold=True, color=BLUE)
rqs = [
    "RQ1.  Query Rewrite를 적용하면 RAG 검색 단계의 Recall@k 및 MRR이 유의미하게 향상되는가?",
    "RQ2.  BM25와 벡터 검색을 결합한 Hybrid Search는 단일 검색 방식 대비 성능이 향상되는가?",
    "RQ3.  두 기법을 함께 적용했을 때 생성 품질(EM, F1)이 더욱 개선되는가?",
]
rqy = Inches(5.9)
for rq in rqs:
    add_textbox(sld, Inches(0.55), rqy, Inches(12.3), Inches(0.32),
                rq, font_size=Pt(12), color=DARK)
    rqy += Inches(0.35)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — 문제 정의 (2): 연구 목표 및 범위
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "문제 정의", "연구 목표 및 범위", "04 / 14")

# 목표 박스
bullet_box(sld, Inches(0.35), Inches(1.35), Inches(12.6), Inches(2.2),
           "🎯  연구 목표",
           ["① Query Rewrite 모듈을 통해 사용자 질의를 재구성, 검색 단계의 질을 향상시킨다.",
            "② BM25 + 벡터 검색 기반 Hybrid Search를 구현하여 검색 커버리지와 정밀도를 높인다.",
            "③ 기본 RAG / Query Rewrite 적용 / Hybrid Search 적용 / 두 기법 결합 의 4가지 모델을 비교 평가한다.",
            "④ 검색 지표(Recall@k, MRR)와 생성 지표(EM, F1)를 함께 측정하여 최적 구성을 도출한다."],
           title_color=BLUE, item_size=Pt(13))

# 연구 범위 (2열)
lw = Inches(6.1)
rw = Inches(6.5)
bullet_box(sld, Inches(0.35), Inches(3.8), lw, Inches(2.55),
           "📌  연구 포함 범위 (In-scope)",
           ["• 기업 내부 문서 기반 한국어 QA 챗봇",
            "• Query Rewrite (LLM 기반 쿼리 재작성)",
            "• Hybrid Search (BM25 + Dense Vector, Elasticsearch)",
            "• 멀티링구얼 임베딩 (multilingual-e5-large)",
            "• 4가지 구성 간 정량적 비교 실험"],
           title_color=GREEN, border_color=GREEN, item_size=Pt(12.5))

bullet_box(sld, Inches(6.6), Inches(3.8), rw, Inches(2.55),
           "🚫  연구 제외 범위 (Out-of-scope)",
           ["• 실시간 웹 크롤링 및 외부 지식 연동",
            "• 멀티모달 (이미지·표 등) 처리",
            "• 파인튜닝(Fine-tuning) 기반 LLM 최적화",
            "• 다국어 동시 지원 (한국어에 집중)",
            "• 대규모 프로덕션 배포 환경 최적화"],
           title_color=RED, border_color=RED, item_size=Pt(12.5))


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — 제안 기법 (1): 전체 시스템 아키텍처
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "제안 기법", "전체 시스템 아키텍처", "05 / 14")

# ── 파이프라인 플로우 다이어그램 ──
BOX_W  = Inches(1.65)
BOX_H  = Inches(0.72)
ARR_W  = Inches(0.35)
flow_y = Inches(1.55)
steps = [
    ("사용자 질의", BLUE),
    ("Query\nRewrite", RGBColor(0x6B, 0x21, 0xA8)),
    ("Query\nEmbedding", RGBColor(0x05, 0x96, 0x69)),
    ("Hybrid\nSearch", RGBColor(0xD9, 0x7A, 0x06)),
    ("Context\n결합", RGBColor(0x0E, 0x6E, 0x7E)),
    ("LLM\n생성", RGBColor(0xDC, 0x26, 0x26)),
    ("최종 답변", NAVY),
]
total_w = len(steps) * BOX_W + (len(steps) - 1) * ARR_W
start_x = (SLD_W - total_w) / 2
cx = start_x
for i, (label, color) in enumerate(steps):
    add_rect(sld, cx, flow_y, BOX_W, BOX_H, fill_rgb=color)
    add_textbox(sld, cx + Inches(0.05), flow_y + Inches(0.1),
                BOX_W - Inches(0.1), BOX_H - Inches(0.15),
                label, font_size=Pt(11.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cx += BOX_W
    if i < len(steps) - 1:
        add_textbox(sld, cx + Inches(0.04), flow_y + Inches(0.18),
                    ARR_W - Inches(0.08), BOX_H - Inches(0.3),
                    "▶", font_size=Pt(16), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        cx += ARR_W

# 컴포넌트 설명 (3열)
comp_y = Inches(2.65)
comp_h = Inches(3.85)
comps = [
    ("🗄  데이터 레이어",
     BLUE,
     ["MySQL\n  └ 원본 문서 저장 (doc_id, title, content, source)",
      "Elasticsearch\n  └ Dense Vector Index (1024-dim, cosine)",
      "multilingual-e5-large\n  └ 청크(512 token / stride 12) 임베딩"]),
    ("⚙️  검색 레이어",
     RGBColor(0x6B, 0x21, 0xA8),
     ["Query Rewrite Module\n  └ GPT-4o-mini 로 질의 재작성",
      "BM25 (희소 검색)\n  └ 키워드 기반 TF-IDF 스코어링",
      "Vector Search (밀집 검색)\n  └ 코사인 유사도 Top-k 검색",
      "Hybrid Score Fusion\n  └ α·BM25 + (1-α)·Vector 결합"]),
    ("🤖  생성 레이어",
     RGBColor(0x05, 0x96, 0x69),
     ["Context 조합\n  └ Top-k 문서를 Prompt Template에 삽입",
      "GPT-4o-mini 생성\n  └ Temperature 0.3, Max Tokens 512",
      "Streamlit UI\n  └ 챗봇 인터페이스 + ES 대시보드"]),
]
cw = Inches(4.1)
cx = Inches(0.35)
for ctitle, ccolor, citems in comps:
    bullet_box(sld, cx, comp_y, cw, comp_h, ctitle, citems,
               title_color=ccolor, border_color=ccolor, item_size=Pt(11.5))
    cx += cw + Inches(0.26)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — 제안 기법 (2): Query Rewrite 모듈
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "제안 기법", "Query Rewrite 모듈", "06 / 14")

# 왼쪽: 개념 + 알고리즘
lw = Inches(6.4)
bullet_box(sld, Inches(0.35), Inches(1.35), lw, Inches(2.2),
           "💡  Query Rewrite 개념",
           ["• Ma et al. (2023) 제안: LLM을 사용해 사용자 질의를\n"
            "  검색에 최적화된 형태로 재작성",
            "• 구어체·약어·오탈자·문맥 생략 → 명시적 검색 쿼리로 변환",
            "• 검색 단계 진입 전 적용 → 임베딩 품질 향상"],
           title_color=RGBColor(0x6B, 0x21, 0xA8),
           border_color=RGBColor(0x6B, 0x21, 0xA8), item_size=Pt(13))

# 알고리즘 박스
add_rect(sld, Inches(0.35), Inches(3.72), lw, Inches(2.85),
         fill_rgb=RGBColor(0x1E, 0x1E, 0x2E), line_rgb=RGBColor(0x6B, 0x21, 0xA8), line_width=Pt(1.5))
add_textbox(sld, Inches(0.5), Inches(3.8), Inches(5.8), Inches(0.35),
            "Algorithm 1: Query Rewrite", font_size=Pt(12), bold=True,
            color=RGBColor(0xBB, 0x86, 0xFC))
algo_lines = [
    "Input : user_query q",
    "Output: rewritten query q′",
    "",
    "1.  system_prompt ← '쿼리를 검색에 최적화된 형태로 재작성하라'",
    "2.  messages ← [(system, system_prompt), (user, q)]",
    "3.  q′ ← LLM.generate(messages, model='gpt-4o-mini', T=0.3)",
    "4.  return q′",
]
ay = Inches(4.18)
for line in algo_lines:
    add_textbox(sld, Inches(0.5), ay, Inches(5.8), Inches(0.28),
                line, font_size=Pt(11), color=RGBColor(0xA8, 0xFF, 0xC2),
                italic=(not line.startswith("Input") and not line.startswith("Output")))
    ay += Inches(0.28)

# 오른쪽: 예시
rw = Inches(6.4)
rx = Inches(6.9)
add_rect(sld, rx, Inches(1.35), rw, Inches(1.9),
         fill_rgb=WHITE, line_rgb=RGBColor(0x6B, 0x21, 0xA8), line_width=Pt(1.5))
add_rect(sld, rx, Inches(1.35), rw, Inches(0.42),
         fill_rgb=RGBColor(0x6B, 0x21, 0xA8))
add_textbox(sld, rx + Inches(0.15), Inches(1.37), rw - Inches(0.3), Inches(0.38),
            "📝  변환 예시", font_size=Pt(14), bold=True, color=WHITE)
examples = [
    ("입력 질의", "\"코리아테크 23년 상반기 매출 얼마야?\"",
     RGBColor(0xDC, 0x26, 0x26), Inches(1.82)),
    ("재작성 결과", "\"코리아테크 주식회사의 2023년 상반기 매출액은 얼마입니까?\"",
     GREEN, Inches(2.52)),
]
for label, text, color, ey in examples:
    add_textbox(sld, rx + Inches(0.15), ey, Inches(1.3), Inches(0.32),
                label + " :", font_size=Pt(11.5), bold=True, color=color)
    add_textbox(sld, rx + Inches(1.5), ey, rw - Inches(1.6), Inches(0.38),
                text, font_size=Pt(11.5), color=DARK)

# 효과 카드
bullet_box(sld, rx, Inches(3.42), rw, Inches(3.15),
           "📈  예상 효과",
           ["• 검색 쿼리의 의미적 명확성 향상\n"
            "  → 임베딩 공간에서 관련 문서와의 유사도 상승",
            "• 약어·오탈자 정규화\n"
            "  → BM25 키워드 매칭 정확도 향상",
            "• Recall@k 및 MRR 지표 개선 기대\n"
            "  (기본 RAG 대비 성능 향상 검증 예정)"],
           title_color=RGBColor(0x6B, 0x21, 0xA8),
           border_color=RGBColor(0x6B, 0x21, 0xA8), item_size=Pt(12.5))


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — 제안 기법 (3): Hybrid Search
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "제안 기법", "Hybrid Search (BM25 + Vector Search)", "07 / 14")

# 상단 비교표
add_rect(sld, Inches(0.35), Inches(1.35), Inches(12.6), Inches(2.25),
         fill_rgb=WHITE, line_rgb=BLUE, line_width=Pt(1))
# 헤더 행
heads = ["", "BM25 (희소 검색)", "Vector Search (밀집 검색)", "Hybrid Search (제안)"]
col_ws = [Inches(1.8), Inches(3.3), Inches(3.7), Inches(3.8)]
hx = Inches(0.35)
add_rect(sld, hx, Inches(1.35), SLD_W - Inches(0.7), Inches(0.42), fill_rgb=NAVY)
for head, cw in zip(heads, col_ws):
    add_textbox(sld, hx + Inches(0.05), Inches(1.37), cw - Inches(0.1), Inches(0.38),
                head, font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hx += cw

rows = [
    ("원리",     "TF-IDF 기반 토큰 빈도 점수",
                 "임베딩 공간 코사인 유사도",
                 "α·BM25_score + (1-α)·Vector_score"),
    ("장점",     "정확 키워드 매칭 강력",
                 "의미적 유사도 우수",
                 "두 방식의 장점 결합"),
    ("단점",     "의미 파악 한계",
                 "정확 키워드 매칭 취약",
                 "가중치 α 튜닝 필요"),
]
row_colors = [WHITE, RGBColor(0xF8, 0xFA, 0xFF), WHITE]
ry = Inches(1.77)
for (rh, bm25, vec, hyb), rbg in zip(rows, row_colors):
    add_rect(sld, Inches(0.35), ry, SLD_W - Inches(0.7), Inches(0.48), fill_rgb=rbg)
    vals = [rh, bm25, vec, hyb]
    rx_ = Inches(0.35)
    for val, cw in zip(vals, col_ws):
        bold = (val == rh)
        clr = BLUE if bold else DARK
        add_textbox(sld, rx_ + Inches(0.08), ry + Inches(0.06),
                    cw - Inches(0.15), Inches(0.38),
                    val, font_size=Pt(11.5), bold=bold, color=clr)
        rx_ += cw
    ry += Inches(0.48)

# 수식 박스
add_rect(sld, Inches(0.35), Inches(3.77), Inches(12.6), Inches(0.82),
         fill_rgb=RGBColor(0x1E, 0x1E, 0x2E))
add_textbox(sld, Inches(0.5), Inches(3.83), Inches(12.3), Inches(0.32),
            "Hybrid Score(d) = α × BM25(q, d)  +  (1-α) × CosineSim(embed(q), embed(d))",
            font_size=Pt(16), bold=True, color=RGBColor(0xA8, 0xFF, 0xC2),
            align=PP_ALIGN.CENTER)
add_textbox(sld, Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.32),
            "단, α ∈ [0, 1] 은 BM25와 Vector Score 간 가중치 파라미터 (기본값 α=0.5)",
            font_size=Pt(11.5), color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER,
            italic=True)

# 구현 세부
lw = Inches(6.1)
rw = Inches(6.5)
bullet_box(sld, Inches(0.35), Inches(4.72), lw, Inches(2.32),
           "🔧  구현 세부 (BM25)",
           ["• Elasticsearch 내장 BM25 스코어러 활용",
            "• 한국어 형태소 분석기 (nori analyzer) 적용\n"
            "  → 어절 분리 및 조사 제거로 정확도 향상",
            "• 필드: content, title 대상 멀티 필드 검색"],
           title_color=RGBColor(0xD9, 0x7A, 0x06),
           border_color=RGBColor(0xD9, 0x7A, 0x06), item_size=Pt(12))

bullet_box(sld, Inches(6.6), Inches(4.72), rw, Inches(2.32),
           "🔧  구현 세부 (Vector)",
           ["• multilingual-e5-large 모델 (1024-dim 벡터)",
            "• Elasticsearch kNN 검색으로 Top-k 후보 추출",
            "• 청크 단위(512 token) 검색 후\n"
            "  원본 문서 단위로 Re-ranking"],
           title_color=BLUE, border_color=BLUE, item_size=Pt(12))


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — 제안 기법 (4): RAG 파이프라인 상세 흐름
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "제안 기법", "RAG 파이프라인 상세 흐름 (Indexing & Retrieval)", "08 / 14")

# ── Indexing 파이프라인 ──
add_rect(sld, Inches(0.35), Inches(1.35), Inches(12.6), Inches(0.38),
         fill_rgb=GREEN)
add_textbox(sld, Inches(0.5), Inches(1.37), Inches(12.3), Inches(0.35),
            "① 문서 색인 (Indexing Pipeline)",
            font_size=Pt(13), bold=True, color=WHITE)

idx_steps = [
    ("원본\n문서", "기업 내부\n문서 (.txt,\n.pdf 등)"),
    ("청크\n분할", "512 token,\nstride=12\noverlap"),
    ("임베딩\n생성", "multilingual\n-e5-large\n(1024-dim)"),
    ("MySQL\n저장", "doc_id,\ntitle, content,\nis_embedded"),
    ("ES\n색인", "dense_vector\n+ BM25\n색인 동시 저장"),
]
ix = Inches(0.5)
for i, (label, desc) in enumerate(idx_steps):
    add_rect(sld, ix, Inches(1.85), Inches(2.2), Inches(1.2), fill_rgb=GREEN)
    add_textbox(sld, ix + Inches(0.05), Inches(1.9),
                Inches(2.1), Inches(0.4),
                label, font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sld, ix + Inches(0.05), Inches(2.3),
                Inches(2.1), Inches(0.7),
                desc, font_size=Pt(10.5), color=RGBColor(0xD4, 0xF5, 0xD4),
                align=PP_ALIGN.CENTER)
    ix += Inches(2.2)
    if i < len(idx_steps) - 1:
        add_textbox(sld, ix, Inches(2.2), Inches(0.25), Inches(0.42),
                    "▶", font_size=Pt(14), bold=True, color=ACCENT)
        ix += Inches(0.25)

# ── Retrieval 파이프라인 ──
add_rect(sld, Inches(0.35), Inches(3.22), Inches(12.6), Inches(0.38),
         fill_rgb=BLUE)
add_textbox(sld, Inches(0.5), Inches(3.24), Inches(12.3), Inches(0.35),
            "② 검색 및 생성 (Retrieval & Generation Pipeline)",
            font_size=Pt(13), bold=True, color=WHITE)

ret_steps = [
    ("사용자\n질의 입력", "Streamlit\nUI"),
    ("Query\nRewrite", "GPT-4o-mini\n로 쿼리 재작성"),
    ("쿼리\n임베딩", "e5-large\n1024-dim"),
    ("Hybrid\nSearch", "BM25 + Vector\nα=0.5 가중 합산"),
    ("Top-k\n문서 선택", "k=3\n유사도 점수 정렬"),
    ("Prompt\n구성", "Template에\nContext 삽입"),
    ("LLM\n생성", "GPT-4o-mini\nT=0.3"),
]
rx2 = Inches(0.35)
rbox_w = (SLD_W - Inches(0.7)) / len(ret_steps) - Inches(0.05)
for i, (label, desc) in enumerate(ret_steps):
    add_rect(sld, rx2, Inches(3.72), rbox_w, Inches(1.25), fill_rgb=BLUE)
    add_textbox(sld, rx2 + Inches(0.04), Inches(3.77),
                rbox_w - Inches(0.08), Inches(0.4),
                label, font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sld, rx2 + Inches(0.04), Inches(4.18),
                rbox_w - Inches(0.08), Inches(0.72),
                desc, font_size=Pt(9.5), color=RGBColor(0xB0, 0xD4, 0xFF),
                align=PP_ALIGN.CENTER)
    rx2 += rbox_w + Inches(0.05)

# 핵심 포인트
bullet_box(sld, Inches(0.35), Inches(5.18), Inches(12.6), Inches(2.05),
           "⭐  파이프라인 핵심 설계 포인트",
           ["• Query Rewrite는 임베딩 및 BM25 검색 이전에 적용 → 검색 전 단계에서 질의 품질 선제적 향상",
            "• Hybrid Search Score Fusion 후 Top-k를 선택 → 단일 방식 대비 다양한 관련 문서 포함 가능",
            "• Prompt Template에 검색 문서를 순위별로 삽입 → LLM이 가장 관련성 높은 컨텍스트부터 처리"],
           title_color=ACCENT, border_color=ACCENT, item_size=Pt(12.5))


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — 제안 기법 (5): 구현 현황
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "제안 기법", "구현 현황 및 기술 스택", "09 / 14")

# 구현 현황 테이블
add_rect(sld, Inches(0.35), Inches(1.35), Inches(12.6), Inches(3.35),
         fill_rgb=WHITE, line_rgb=BLUE, line_width=Pt(1))

headers = ["모듈", "구현 파일", "기술 스택", "구현 상태"]
col_ws2 = [Inches(2.2), Inches(3.0), Inches(4.2), Inches(3.2)]
add_rect(sld, Inches(0.35), Inches(1.35), SLD_W - Inches(0.7), Inches(0.42), fill_rgb=NAVY)
hx2 = Inches(0.35)
for head, cw in zip(headers, col_ws2):
    add_textbox(sld, hx2 + Inches(0.08), Inches(1.37), cw - Inches(0.15), Inches(0.38),
                head, font_size=Pt(12.5), bold=True, color=WHITE)
    hx2 += cw

impl_rows = [
    ("문서 임베딩", "embedder/\nchunk_embedder.py",
     "multilingual-e5-large, HuggingFace\nTransformers, Sentence-Transformers", "✅ 완료"),
    ("Vector 검색", "retriever/\nelastic_search.py",
     "Elasticsearch 8.x, kNN search\ncosine similarity", "✅ 완료"),
    ("LLM 생성", "generator/\ngpt_generator.py",
     "OpenAI API, GPT-4o-mini\nPrompt Template", "✅ 완료"),
    ("DB 연동", "db/mysql_connector.py",
     "MySQL, python-mysql-connector\nUUID 기반 문서 관리", "✅ 완료"),
    ("UI / 대시보드", "app.py",
     "Streamlit, ES 대시보드\n인덱스 상태·문서 조회", "✅ 완료"),
    ("Query Rewrite", "retriever/ (예정)",
     "GPT-4o-mini 쿼리 재작성 모듈", "🔧 구현 예정"),
    ("Hybrid Search", "retriever/ (예정)",
     "BM25 + Vector Score Fusion (α=0.5)", "🔧 구현 예정"),
    ("평가 시스템", "eval/ (예정)",
     "Recall@k, MRR, EM, F1 측정 스크립트", "🔧 구현 예정"),
]
row_y2 = Inches(1.77)
for i, (mod, file_, stack, status) in enumerate(impl_rows):
    rbg = WHITE if i % 2 == 0 else RGBColor(0xF3, 0xF4, 0xF6)
    status_color = GREEN if "✅" in status else ACCENT
    add_rect(sld, Inches(0.35), row_y2, SLD_W - Inches(0.7), Inches(0.41), fill_rgb=rbg)
    vals = [mod, file_, stack, status]
    rx3 = Inches(0.35)
    for j, (val, cw) in enumerate(zip(vals, col_ws2)):
        clr = status_color if j == 3 else DARK
        bld = (j == 3)
        add_textbox(sld, rx3 + Inches(0.08), row_y2 + Inches(0.03),
                    cw - Inches(0.15), Inches(0.36),
                    val, font_size=Pt(10.5), bold=bld, color=clr)
        rx3 += cw
    row_y2 += Inches(0.41)

# 기술 스택 요약
bullet_box(sld, Inches(0.35), Inches(4.9), Inches(12.6), Inches(1.82),
           "🛠  전체 기술 스택 요약",
           ["Frontend: Streamlit (Python 기반 웹 앱)  |  "
            "Backend: Python 3.10+, FastAPI (예정)",
            "검색 엔진: Elasticsearch 8.x (kNN + BM25)  |  "
            "DB: MySQL 8.0",
            "임베딩: multilingual-e5-large (HuggingFace)  |  "
            "LLM: OpenAI GPT-4o-mini  |  "
            "환경: AWS EC2 (Ubuntu)"],
           title_color=NAVY, border_color=NAVY, item_size=Pt(13))


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — 검증 (1): 실험 설계
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "제안 기법 검증", "실험 설계", "10 / 14")

# 4가지 모델 비교 설계
add_rect(sld, Inches(0.35), Inches(1.35), Inches(12.6), Inches(0.38),
         fill_rgb=BLUE)
add_textbox(sld, Inches(0.5), Inches(1.37), Inches(12.3), Inches(0.35),
            "비교 대상 모델 (Baselines vs. Proposed)",
            font_size=Pt(13), bold=True, color=WHITE)

model_info = [
    ("M1", "Basic RAG",       "Vector Search Only\n(Baseline)",
     GRAY, "기본 벡터 검색만 사용"),
    ("M2", "+ Query Rewrite", "Vector Search\n+ Query Rewrite",
     RGBColor(0x6B, 0x21, 0xA8), "쿼리 재작성 단독 효과"),
    ("M3", "+ Hybrid Search", "BM25 + Vector\n(No Rewrite)",
     RGBColor(0xD9, 0x7A, 0x06), "하이브리드 검색 단독 효과"),
    ("M4", "Full Proposed",   "Query Rewrite\n+ Hybrid Search",
     GREEN, "두 기법 결합 (최종 제안)"),
]
mx = Inches(0.35)
mw = Inches(3.1)
for mid, mname, mdesc, mcolor, mrole in model_info:
    add_rect(sld, mx, Inches(1.8), mw, Inches(1.35), fill_rgb=mcolor)
    add_textbox(sld, mx + Inches(0.08), Inches(1.84), mw - Inches(0.15), Inches(0.4),
                f"{mid}  {mname}", font_size=Pt(13), bold=True, color=WHITE)
    add_textbox(sld, mx + Inches(0.08), Inches(2.24), mw - Inches(0.15), Inches(0.65),
                mdesc, font_size=Pt(11), color=RGBColor(0xE0, 0xF0, 0xFF))
    add_rect(sld, mx, Inches(3.15), mw, Inches(0.38),
             fill_rgb=RGBColor(0xF0, 0xF4, 0xFF), line_rgb=mcolor, line_width=Pt(1))
    add_textbox(sld, mx + Inches(0.08), Inches(3.18), mw - Inches(0.15), Inches(0.32),
                mrole, font_size=Pt(10.5), color=mcolor, italic=True)
    mx += mw + Inches(0.17)

# 데이터셋 + 환경
lw2 = Inches(6.1)
rw2 = Inches(6.5)
bullet_box(sld, Inches(0.35), Inches(3.75), lw2, Inches(2.35),
           "📊  데이터셋 구성 (계획)",
           ["• 합성 기업 문서: GPT-4o-mini로 생성\n"
            "  (한국어 중견기업 설정, JSON 형식)",
            "• 문서 수: ~50개 기업 문서 (각 500자 이상)",
            "• 청크 단위 분할 후 총 ~200-300 청크",
            "• 평가용 QA 페어: 각 문서에서 질문-답변\n"
            "  10~20쌍 생성 (총 100~200 QA)"],
           title_color=RGBColor(0x0E, 0x6E, 0x7E),
           border_color=RGBColor(0x0E, 0x6E, 0x7E), item_size=Pt(12))

bullet_box(sld, Inches(6.6), Inches(3.75), rw2, Inches(2.35),
           "🖥  실험 환경",
           ["• OS: Ubuntu 22.04 (AWS EC2)",
            "• Elasticsearch: 8.x (kNN + BM25)",
            "• 임베딩 모델: multilingual-e5-large (CPU/GPU)",
            "• LLM: OpenAI GPT-4o-mini (API 호출)",
            "• 프레임워크: Python 3.10, Streamlit"],
           title_color=NAVY, border_color=NAVY, item_size=Pt(12))


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — 검증 (2): 평가 지표 및 예상 결과
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "제안 기법 검증", "평가 지표 및 예상 실험 결과", "11 / 14")

# 지표 설명
add_rect(sld, Inches(0.35), Inches(1.35), Inches(12.6), Inches(0.38),
         fill_rgb=BLUE)
add_textbox(sld, Inches(0.5), Inches(1.37), Inches(12.3), Inches(0.35),
            "평가 지표", font_size=Pt(13), bold=True, color=WHITE)

metrics = [
    ("Recall@k",  "검색 단계",
     "관련 문서가 Top-k 안에 포함된 비율\n높을수록 검색 커버리지 우수",
     BLUE),
    ("MRR",       "검색 단계",
     "Mean Reciprocal Rank\n첫 번째 관련 문서의 순위 역수 평균",
     RGBColor(0x6B, 0x21, 0xA8)),
    ("Exact Match","생성 단계",
     "생성 답변이 정답과 완전 일치하는 비율",
     RGBColor(0xD9, 0x7A, 0x06)),
    ("F1 Score",  "생성 단계",
     "토큰 수준 정밀도-재현율의 조화 평균\n부분 일치 성능 측정",
     GREEN),
]
mx2 = Inches(0.35)
mw2 = Inches(3.1)
for mname, mstage, mdesc, mcolor in metrics:
    add_rect(sld, mx2, Inches(1.82), mw2, Inches(1.55),
             fill_rgb=WHITE, line_rgb=mcolor, line_width=Pt(2))
    add_rect(sld, mx2, Inches(1.82), mw2, Inches(0.42), fill_rgb=mcolor)
    add_textbox(sld, mx2 + Inches(0.1), Inches(1.84),
                mw2 - Inches(0.2), Inches(0.38),
                f"{mname}  [{mstage}]", font_size=Pt(12.5), bold=True, color=WHITE)
    add_textbox(sld, mx2 + Inches(0.1), Inches(2.3),
                mw2 - Inches(0.2), Inches(1.0),
                mdesc, font_size=Pt(11.5), color=DARK)
    mx2 += mw2 + Inches(0.17)

# 예상 결과 테이블
add_rect(sld, Inches(0.35), Inches(3.55), Inches(12.6), Inches(0.38),
         fill_rgb=NAVY)
add_textbox(sld, Inches(0.5), Inches(3.57), Inches(12.3), Inches(0.35),
            "예상 실험 결과 (가설 기반 추정치)",
            font_size=Pt(13), bold=True, color=WHITE)

result_cols = ["모델", "Recall@3", "MRR", "Exact Match", "F1 Score", "비고"]
result_ws = [Inches(2.4), Inches(1.85), Inches(1.85), Inches(2.1), Inches(2.1), Inches(2.3)]

add_rect(sld, Inches(0.35), Inches(3.93), SLD_W - Inches(0.7), Inches(0.38),
         fill_rgb=RGBColor(0x1A, 0x56, 0xDB))
hrx = Inches(0.35)
for hh, hw in zip(result_cols, result_ws):
    add_textbox(sld, hrx + Inches(0.05), Inches(3.95), hw - Inches(0.1), Inches(0.33),
                hh, font_size=Pt(11.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hrx += hw

result_data = [
    ("M1  Basic RAG",         "0.62", "0.55", "0.41", "0.48", "Baseline"),
    ("M2  + Query Rewrite",   "0.71", "0.63", "0.47", "0.54", "↑ QR 효과"),
    ("M3  + Hybrid Search",   "0.75", "0.68", "0.50", "0.57", "↑ HS 효과"),
    ("M4  Full (QR+HS) ★",    "0.83", "0.76", "0.57", "0.65", "최종 제안"),
]
rry = Inches(4.31)
row_bgs = [WHITE, RGBColor(0xF8, 0xFA, 0xFF), WHITE, RGBColor(0xF0, 0xF9, 0xF0)]
for i, (row_data, rbg) in enumerate(zip(result_data, row_bgs)):
    add_rect(sld, Inches(0.35), rry, SLD_W - Inches(0.7), Inches(0.42), fill_rgb=rbg)
    rrx = Inches(0.35)
    for j, (val, hw) in enumerate(zip(row_data, result_ws)):
        highlight = (i == 3)
        clr = GREEN if (highlight and j > 0) else (NAVY if j == 0 and highlight else DARK)
        bld = highlight or j == 0
        add_textbox(sld, rrx + Inches(0.05), rry + Inches(0.05),
                    hw - Inches(0.1), Inches(0.33),
                    val, font_size=Pt(11.5), bold=bld, color=clr,
                    align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
        rrx += hw
    rry += Inches(0.42)

add_textbox(sld, Inches(0.35), Inches(6.7), Inches(12.6), Inches(0.28),
            "※ 수치는 실험 설계 단계의 가설적 추정치이며, 실제 구현 후 갱신 예정",
            font_size=Pt(10), color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — 검증 (3): 실험 프로세스 및 분석 방법
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "제안 기법 검증", "실험 프로세스 및 분석 방법", "12 / 14")

# 실험 절차
bullet_box(sld, Inches(0.35), Inches(1.35), Inches(12.6), Inches(2.5),
           "🔬  실험 절차",
           ["① [데이터 준비]  GPT-4o-mini로 합성 기업 문서 50개 생성 → 청크 분할 → ES 색인",
            "② [QA 구성]  각 문서에서 질문-정답 페어 자동 생성 (GPT 활용) → 수동 검토",
            "③ [모델 실행]  M1→M4 순서로 동일 QA셋에 대해 각 모델 검색 결과 및 생성 답변 저장",
            "④ [지표 계산]  Recall@k, MRR (검색 단계) / EM, F1 (생성 단계) 자동 계산 스크립트 실행",
            "⑤ [통계 분석]  4개 모델 간 성능 차이의 통계적 유의성 검증 (paired t-test)"],
           title_color=BLUE, item_size=Pt(13))

# 분석 포인트
lw3 = Inches(6.1)
rw3 = Inches(6.5)
bullet_box(sld, Inches(0.35), Inches(4.05), lw3, Inches(2.7),
           "📐  분석 관점",
           ["• Query Rewrite 단독 효과\n"
            "  → M1 vs. M2 비교",
            "• Hybrid Search 단독 효과\n"
            "  → M1 vs. M3 비교",
            "• 두 기법 결합 시너지\n"
            "  → M4 vs. M2·M3 비교",
            "• 검색 지표와 생성 지표의\n"
            "  상관관계 분석"],
           title_color=RGBColor(0x0E, 0x6E, 0x7E),
           border_color=RGBColor(0x0E, 0x6E, 0x7E), item_size=Pt(12.5))

bullet_box(sld, Inches(6.6), Inches(4.05), rw3, Inches(2.7),
           "⚠️  실험 제한사항 및 대응",
           ["• 합성 데이터 사용으로 인한 실제 도메인\n"
            "  일반화 한계 → 가능 시 실제 기업 문서\n"
            "  일부 포함 예정",
            "• QA 자동 생성 품질 불확실\n"
            "  → 수동 검토 및 필터링",
            "• α 파라미터 민감도\n"
            "  → α ∈ {0.3, 0.5, 0.7} Grid Search"],
           title_color=RED, border_color=RED, item_size=Pt(12.5))


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — 결론 및 공헌점
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "결론 및 공헌점", "", "13 / 14")

# 연구 요약
bullet_box(sld, Inches(0.35), Inches(1.35), Inches(12.6), Inches(1.85),
           "📋  연구 요약",
           ["• 기업 환경에서 내부 문서 기반 QA 챗봇 구축 시 기존 단순 RAG의 검색 품질 한계를 식별",
            "• Query Rewrite + Hybrid Search (BM25 + Vector) 결합 기법을 적용한 확장 RAG 아키텍처 제안",
            "• 4가지 구성 모델의 정량적 비교 실험을 통해 각 기법의 단독·결합 효과를 체계적으로 분석"],
           title_color=NAVY, item_size=Pt(13))

# 공헌점 4개
contrib_y = Inches(3.38)
contribs = [
    ("🎯", "실용적 아키텍처 제안",
     "LLM 기반 Query Rewrite와 Elasticsearch Hybrid Search를\n실제 구현 가능한 수준으로 통합 설계"),
    ("📊", "체계적 비교 평가",
     "검색 단계(Recall@k, MRR)와 생성 단계(EM, F1)를\n함께 측정하는 통합 평가 프레임워크 수립"),
    ("🔧", "오픈소스 구현체 제공",
     "Streamlit 기반 챗봇 UI + Elasticsearch 대시보드\n통합 시스템을 공개 코드로 제공 예정"),
    ("🌐", "한국어 특화 적용",
     "multilingual-e5-large + nori 형태소 분석기를\n적용한 한국어 기업 환경 최적화"),
]
cx2 = Inches(0.35)
cw2 = Inches(3.1)
for icon, ctitle, cdesc in contribs:
    add_rect(sld, cx2, contrib_y, cw2, Inches(2.35),
             fill_rgb=WHITE, line_rgb=BLUE, line_width=Pt(1.5))
    add_rect(sld, cx2, contrib_y, cw2, Inches(0.5), fill_rgb=BLUE)
    add_textbox(sld, cx2 + Inches(0.1), contrib_y + Inches(0.05),
                cw2 - Inches(0.2), Inches(0.42),
                f"{icon}  {ctitle}", font_size=Pt(12.5), bold=True, color=WHITE)
    add_textbox(sld, cx2 + Inches(0.1), contrib_y + Inches(0.58),
                cw2 - Inches(0.2), Inches(1.65),
                cdesc, font_size=Pt(12), color=DARK)
    cx2 += cw2 + Inches(0.17)

# 향후 계획
add_rect(sld, Inches(0.35), Inches(5.93), Inches(12.6), Inches(0.7),
         fill_rgb=RGBColor(0xFE, 0xF9, 0xEE), line_rgb=ACCENT, line_width=Pt(1.5))
add_textbox(sld, Inches(0.5), Inches(5.98), Inches(1.5), Inches(0.4),
            "향후 계획", font_size=Pt(13), bold=True, color=ACCENT)
add_textbox(sld, Inches(2.1), Inches(5.98), Inches(10.7), Inches(0.55),
            "Query Rewrite 모듈 구현 → Hybrid Search 구현 → 평가 데이터셋 구성 → "
            "4-모델 비교 실험 수행 → 논문 작성 (2026년 하반기 목표)",
            font_size=Pt(12.5), color=RGBColor(0x92, 0x40, 0x00))


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — 참고문헌 + Q&A
# ═══════════════════════════════════════════════════════════════
sld = prs.slides.add_slide(BLANK)
slide_header(sld, "참고문헌", "", "14 / 14")

refs = [
    "[1] Lewis, P. et al. (2020). Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks. "
    "NeurIPS 2020.",
    "[2] Karpukhin, V. et al. (2020). Dense Passage Retrieval for Open-Domain Question Answering. "
    "EMNLP 2020.",
    "[3] Gao, Y. et al. (2024). Retrieval-Augmented Generation for Large Language Models: A Survey. "
    "arXiv:2312.10997.",
    "[4] Ram, O. et al. (2023). In-Context Retrieval-Augmented Language Models. "
    "TACL 2023.",
    "[5] Gao, L. et al. (2024). RAGBench: Explainable Benchmark for Retrieval-Augmented Generation Systems. "
    "arXiv:2407.11005.",
    "[6] Ma, X. et al. (2023). Query Rewriting for Retrieval-Augmented Large Language Models. "
    "EMNLP 2023.",
    "[7] Chen, J. et al. (2022). Hybrid Retrieval with Dense and Sparse Representations. "
    "SIGIR 2022.",
]
ref_y = Inches(1.42)
for ref in refs:
    add_textbox(sld, Inches(0.5), ref_y, Inches(12.3), Inches(0.55),
                ref, font_size=Pt(12), color=DARK)
    ref_y += Inches(0.55)

# Q&A 배너
add_rect(sld, Inches(2.5), Inches(5.5), Inches(8.3), Inches(1.55), fill_rgb=NAVY)
add_textbox(sld, Inches(2.5), Inches(5.62), Inches(8.3), Inches(0.7),
            "감사합니다  |  Q & A",
            font_size=Pt(30), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sld, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.45),
            "권성재  •  인공지능학과  •  지도교수: 노동건 교수님",
            font_size=Pt(13), color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)


# ─── 저장 ──────────────────────────────────────────────────────
out_path = r"c:\rag_chatbot\RAG챗봇_발표자료_권성재.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")